"""
PPT/PDF 文件解析服务
"""
from io import BytesIO


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """从 PPTX 文件提取文本"""
    from pptx import Presentation

    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception:
        raise ValueError(
            "无法解析该 PPT 文件。可能原因：\n"
            "1. 这是旧版 .ppt 格式（97-2003），请用 PowerPoint 另存为 .pptx 格式后重试\n"
            "2. 文件已损坏或不完整"
        )

    texts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))
        if slide_texts:
            texts.append(f"--- 第{i}页 ---\n" + "\n".join(slide_texts))

    return "\n\n".join(texts) if texts else ""


def extract_text_from_docx(file_bytes: bytes) -> str:
    """从 DOCX 文件提取文本"""
    from docx import Document

    try:
        doc = Document(BytesIO(file_bytes))
    except Exception:
        raise ValueError("无法解析该 Word 文件。可能文件已损坏或使用了旧版 .doc 格式")

    texts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)

    # 也提取表格中的文本
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                texts.append(" | ".join(row_texts))

    return "\n".join(texts) if texts else ""


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """从 PDF 文件提取文本"""
    from PyPDF2 import PdfReader

    try:
        reader = PdfReader(BytesIO(file_bytes))
    except Exception:
        raise ValueError("无法解析该 PDF 文件，文件可能已损坏或加密")

    texts: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            texts.append(f"--- 第{i}页 ---\n{text.strip()}")

    return "\n\n".join(texts) if texts else ""


def parse_file(filename: str, file_bytes: bytes) -> tuple[str, str]:
    """解析文件，返回 (文本, 文件类型)"""
    filename_lower = filename.lower()

    if filename_lower.endswith(".pptx"):
        text = extract_text_from_pptx(file_bytes)
        return text, "pptx"

    elif filename_lower.endswith(".ppt"):
        raise ValueError(
            "不支持旧版 .ppt 格式（97-2003 版本）。\n"
            "请用 PowerPoint 打开，点击「文件 → 另存为 → .pptx 格式」，然后重新上传。"
        )

    elif filename_lower.endswith(".docx"):
        text = extract_text_from_docx(file_bytes)
        return text, "docx"

    elif filename_lower.endswith(".doc"):
        raise ValueError(
            "不支持旧版 .doc 格式（97-2003 版本）。\n"
            "请用 Word 打开，点击「文件 → 另存为 → .docx 格式」，然后重新上传。"
        )

    elif filename_lower.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
        return text, "pdf"

    else:
        raise ValueError(f"不支持的文件格式: {filename}。仅支持 PPTX、PDF、DOCX 文件。")
